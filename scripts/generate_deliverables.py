"""QAC — Generate Bloco 3 deliverables from official templates.

Creates:
1. Report (3 pages) from 'Cópia de templateA4.docx'
2. Slides (5 slides + cover) from 'Cópia de [Template Slides] Brazil Quantum Camp.pptx'

Usage:
    python scripts/generate_deliverables.py
"""
from __future__ import annotations

import copy
import os
from pathlib import Path

# ── Report Generation (python-docx) ──────────────────────────────────────────

def generate_report(
    template_path: str,
    output_path: str,
) -> None:
    """Generate 3-page report from templateA4.docx."""
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document(template_path)

    # Clear existing content after template header
    # Keep only the first few paragraphs (header/formatting) and add our content
    # First, understand the template structure
    existing_paras = list(doc.paragraphs)

    # Add a page break to start fresh content area
    # We'll append our content to the document

    def add_heading(doc, text, level=1):
        h = doc.add_heading(text, level=level)
        return h

    def add_paragraph(doc, text, bold=False, size=11):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.font.size = Pt(size)
        run.bold = bold
        return p

    # ── PAGE 1: Problem and Solution ──
    add_heading(doc, "1. Problema e Solução Proposta", level=1)

    add_heading(doc, "1.1 Problema", level=2)
    add_paragraph(doc,
        "A classificação automática de uso do solo agrícola a partir de imagens de satélite "
        "é essencial para agricultura de precisão, monitoramento ambiental e planejamento territorial. "
        "Métodos tradicionais de machine learning (SVM, Random Forest) obtêm alta acurácia, mas "
        "enfrentam limitações de escalabilidade em espaços de alta dimensionalidade. "
        "A computação quântica, via algoritmos variacionais híbridos, oferece uma perspectiva "
        "promissora para explorar correlações não-lineares complexas em dados espectrais."
    )

    add_heading(doc, "1.2 Solução: Variational Quantum Classifier (VQC)", level=2)
    add_paragraph(doc,
        "Propomos um Classificador Quântico Variacional (VQC) para classificação binária "
        "de imagens EuroSAT (AnnualCrop vs SeaLake). O VQC utiliza:\n\n"
        "• ZZFeatureMap: Codificação de features clássicas em estados quânticos via angle encoding "
        "com entangling gates, explorando correlações entre features.\n\n"
        "• RealAmplitudes Ansatz: Circuito parametrizado com rotações Ry e portas CNOT "
        "em 2 repetições, totalizando ~24 parâmetros treináveis.\n\n"
        "• COBYLA Optimizer: Otimizador clássico gradient-free, adequado para "
        "paisagens de custo ruidosas típicas de circuitos quânticos.\n\n"
        "A escolha do VQC (em vez de VQE ou QAOA) justifica-se pela sua "
        "formulação nativa para classificação supervisionada, minimizando diretamente "
        "uma loss function de classificação (Benedetti et al., 2019; Schuld & Petruccione, 2021)."
    )

    add_heading(doc, "1.3 Investigação de Algoritmos", level=2)
    add_paragraph(doc,
        "Foram investigadas 3 abordagens:\n\n"
        "1. QSVM (Quantum Support Vector Machine): Kernel quântico via Fidelity — "
        "alta complexidade computacional O(n²) para construção da kernel matrix.\n\n"
        "2. VQC (Variational Quantum Classifier): Circuito parametrizado com otimização "
        "clássica — abordagem selecionada por aderência ao problema e viabilidade NISQ.\n\n"
        "3. VQE como classificador: Hamiltoniano de Ising condicionado por dados — "
        "rejeitado por inadequação conceitual (VQE resolve eigenvalue problems, não classification)."
    )

    # ── PAGE 2: Data, Experiments, Results ──
    doc.add_page_break()
    add_heading(doc, "2. Dados, Experimentos e Resultados", level=1)

    add_heading(doc, "2.1 Dataset EuroSAT", level=2)
    add_paragraph(doc,
        "O EuroSAT é um dataset de imagens de satélite Sentinel-2 com 10 classes de "
        "uso/cobertura do solo. Para este estudo, selecionamos 2 classes representativas:\n\n"
        "• AnnualCrop (cultivos anuais): 50 amostras\n"
        "• SeaLake (corpos d'água): 50 amostras\n\n"
        "Imagens RGB 64×64 pixels (12.288 features) foram reduzidas via PCA para 8 componentes "
        "principais, retendo aproximadamente 49% da variância total. "
        "As 8 features foram normalizadas para [0, 2π] para encoding quântico."
    )

    add_heading(doc, "2.2 Resultados Experimentais", level=2)
    add_paragraph(doc,
        "A tabela abaixo compara o baseline clássico (SVM kernel RBF) com o "
        "Classificador Quântico Variacional (VQC):\n\n"
        "| Métrica      | SVM    | VQC    |\n"
        "| Accuracy     | 0.95   | 0.55   |\n"
        "| F1 Score     | 0.95   | 0.55   |\n"
        "| Tempo Treino | <1s    | 68s    |\n\n"
        "Nota: Os valores exatos são obtidos pela execução do notebook."
    )

    add_heading(doc, "2.3 Análise", level=2)
    add_paragraph(doc,
        "O SVM supera o VQC neste cenário, confirmando que em espaços de baixa "
        "dimensionalidade (8 features PCA linearmente separáveis), classificadores clássicos "
        "são altamente eficientes. Este resultado é esperado e consistente com a literatura: "
        "a compressão PCA reduz a complexidade geométrica que algoritmos quânticos poderiam explorar. "
        "O sistema operacional e a arquitetura do pipeline estão validados e "
        "prontos para cenários mais complexos."
    )

    # ── PAGE 3: Conclusions and References ──
    doc.add_page_break()
    add_heading(doc, "3. Conclusões, Próximos Passos e Referências", level=1)

    add_heading(doc, "3.1 Conclusões", level=2)
    add_paragraph(doc,
        "O Quantum AgriClassifier (QAC) demonstra um pipeline completo e funcional "
        "de classificação quântica híbrida para dados agrícolas. As principais conclusões são:\n\n"
        "1. O VQC é o algoritmo correto para classificação supervisionada quântica\n"
        "2. Em cenários NISQ com dados linearmente separáveis, baselines clássicos dominam\n"
        "3. A ausência de vantagem quântica é resultado honesto e cientificamente significativo\n"
        "4. A arquitetura SDD + Clean Code garante reprodutibilidade total do experimento"
    )

    add_heading(doc, "3.2 Próximos Passos", level=2)
    add_paragraph(doc,
        "• Data re-uploading: Múltiplas camadas de encoding para aumentar expressividade\n"
        "• EfficientSU2 ansatz: Maior capacidade representativa vs RealAmplitudes\n"
        "• Feature selection estatística: Mutual information em vez de PCA\n"
        "• Noise-aware training: Modelos de ruído realistas via Qiskit Aer\n"
        "• Hardware execution: Validação em IBM Quantum (127+ qubits)"
    )

    add_heading(doc, "3.3 Referências", level=2)
    add_paragraph(doc,
        "[1] Benedetti, M. et al. (2019). Parameterized quantum circuits as machine "
        "learning models. Quantum Science and Technology, 4(4).\n\n"
        "[2] Schuld, M. & Petruccione, F. (2021). Machine Learning with Quantum Computers. "
        "Springer.\n\n"
        "[3] Helber, P. et al. (2019). EuroSAT: A Novel Dataset and Deep Learning Benchmark "
        "for Land Use and Land Cover Classification. IEEE JSTARS.\n\n"
        "[4] Havlíček, V. et al. (2019). Supervised learning with quantum-enhanced feature "
        "spaces. Nature, 567(7747).\n\n"
        "[5] McClean, J. R. et al. (2018). Barren plateaus in quantum neural network "
        "training landscapes. Nature Communications, 9(1).",
        size=10,
    )

    doc.save(output_path)
    print(f"✅ Report saved: {output_path}")


# ── Slides Generation (python-pptx) ──────────────────────────────────────────

def generate_slides(
    template_path: str,
    output_path: str,
) -> None:
    """Generate 5 slides + cover from Brazil Quantum Camp template."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(template_path)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Use the first layout as template for content slides
    # Examine available layouts
    layouts = prs.slide_layouts
    # Use layout index 1 (usually Title + Content) or 0 (Title)
    title_layout = layouts[0] if len(layouts) > 0 else layouts[0]
    content_layout = layouts[1] if len(layouts) > 1 else layouts[0]

    def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None):
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        if bold:
            p.font.bold = True
        if color:
            p.font.color.rgb = color
        return txBox

    def add_slide_with_content(prs, layout, title_text, body_text, bullet_points=None):
        slide = prs.slides.add_slide(layout)
        # Set title if available
        if slide.shapes.title:
            slide.shapes.title.text = title_text

        # Add body content
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(8.4)
        height = Inches(4.5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        if body_text:
            p = tf.paragraphs[0]
            p.text = body_text
            p.font.size = Pt(16)

        if bullet_points:
            for bp in bullet_points:
                p = tf.add_paragraph()
                p.text = bp
                p.font.size = Pt(14)
                p.level = 0

        return slide

    # ── SLIDE 0: COVER ──
    # The template likely has a cover slide already — modify or add
    if len(prs.slides) > 0:
        # Modify existing cover
        cover = prs.slides[0]
        if cover.shapes.title:
            cover.shapes.title.text = "Quantum AgriClassifier (QAC)"
        # Add subtitle
        add_text_box(
            cover, Inches(0.8), Inches(3.5), Inches(8), Inches(1),
            "Classificação Quântica Híbrida de Imagens Agrícolas",
            font_size=20, bold=False,
        )
        add_text_box(
            cover, Inches(0.8), Inches(4.5), Inches(8), Inches(0.8),
            "Leonardo Maximino Bernardo | Equipe Quantum Tech | Bloco 3",
            font_size=14, bold=False,
        )
    else:
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Quantum AgriClassifier (QAC)"

    # ── SLIDE 1: PROBLEMA ──
    add_slide_with_content(
        prs, content_layout,
        "1. Problema",
        "Classificação automática de uso do solo agrícola via imagens de satélite",
        [
            "• Agricultura de precisão requer classificação confiável de cultivos",
            "• Imagens multiespectrais (EuroSAT/Sentinel-2) em alta dimensionalidade",
            "• Métodos clássicos são eficientes, mas enfrentam limites teóricos",
            "• Computação quântica: explorar correlações não-lineares complexas",
            "",
            "Objetivo: Investigar viabilidade do VQC para classificação agrícola",
        ],
    )

    # ── SLIDE 2: SOLUÇÃO ──
    add_slide_with_content(
        prs, content_layout,
        "2. Solução — Variational Quantum Classifier (VQC)",
        "Pipeline: Imagem → PCA(8) → Normalize [0,2π] → VQC → Classe",
        [
            "• ZZFeatureMap: Encoding quântico com entangling (8 qubits)",
            "• RealAmplitudes: Ansatz parametrizado (Ry + CNOT, reps=2)",
            "• COBYLA: Otimizador clássico gradient-free (100 iterações)",
            "• Baseline: SVM (kernel RBF) para comparação",
            "",
            "Fundamentação: Benedetti et al. (2019), Schuld & Petruccione (2021)",
        ],
    )

    # ── SLIDE 3: DADOS/CENÁRIOS ──
    add_slide_with_content(
        prs, content_layout,
        "3. Dados e Cenários",
        "Dataset EuroSAT — Sentinel-2 (RGB 64×64)",
        [
            "• 2 classes: AnnualCrop (cultivos) vs SeaLake (água)",
            "• 100 amostras (50/classe) — cenário Bloco 3",
            "• Features: 12.288 pixels → PCA → 8 componentes (≤10 ✓)",
            "• Variância PCA explicada: ~49%",
            "• Split: 60% treino / 20% val / 20% teste (estratificado)",
            "• Reprodutibilidade: SHA-256 hash + seed fixo (42)",
        ],
    )

    # ── SLIDE 4: RESULTADOS ──
    add_slide_with_content(
        prs, content_layout,
        "4. Resultados e Conclusões",
        "Comparação SVM vs VQC",
        [
            "• SVM: accuracy 95.0%, F1 95.0% (treino <1s)",
            "• VQC: accuracy 55.0%, F1 54.9% (treino ~68s)",
            "• Vantagem quântica: NÃO observada neste cenário",
            "",
            "Análise: PCA lineariza o espaço → SVM domina",
            "O Data-Model Mismatch é independente do domínio dos dados",
            "Resultado negativo honesto = ciência válida",
        ],
    )

    # ── SLIDE 5: LIÇÕES E PRÓXIMOS PASSOS ──
    add_slide_with_content(
        prs, content_layout,
        "5. Lições Aprendidas e Próximos Passos",
        "O que aprendemos e para onde ir",
        [
            "Lições:",
            "  • VQC é o algoritmo correto para classificação quântica",
            "  • NISQ (8 qubits) + PCA = espaço trivial para clássicos",
            "  • Baseline comparativo é obrigatório para validação honesta",
            "",
            "Próximos Passos:",
            "  • Data re-uploading (múltiplas camadas encoding+ansatz)",
            "  • Feature selection estatística (mutual information)",
            "  • Noise-aware training com modelos de ruído",
            "  • Execução em hardware IBM Quantum (127+ qubits)",
        ],
    )

    prs.save(output_path)
    print(f"✅ Slides saved: {output_path}")


# ── Main ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    downloads = Path.home() / "Downloads"
    project = Path.home() / "My_projects" / "Quantum_AgriClassifier_QAC"
    docs_dir = project / "docs"
    docs_dir.mkdir(parents=True, exist_ok=True)

    # Report
    report_template = downloads / "Cópia de templateA4.docx"
    if report_template.exists():
        generate_report(
            str(report_template),
            str(docs_dir / "Entregavel_Bloco_3_Leonardo_Maximino_Bernardo_QAC.docx"),
        )
    else:
        print(f"⚠️ Template not found: {report_template}")

    # Slides
    slides_template = downloads / "Cópia de [Template Slides] Brazil Quantum Camp.pptx"
    if slides_template.exists():
        generate_slides(
            str(slides_template),
            str(docs_dir / "Slides_Bloco_3_Leonardo_Maximino_Bernardo_QAC.pptx"),
        )
    else:
        print(f"⚠️ Template not found: {slides_template}")
